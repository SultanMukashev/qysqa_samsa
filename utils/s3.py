import boto3
import mimetypes
import os
from io import BytesIO

from docx import Document
import PyPDF2
from pptx import Presentation

import openai  # or your LLM library


class S3FileProcessor:
    def __init__(self, bucket_name, aws_access_key='samsas3', aws_secret_key='samsapass', endpoint_url='http://46.101.215.10:9000'):
        if aws_access_key and aws_secret_key:
            self.s3 = boto3.client('s3',
                                   aws_access_key_id=aws_access_key,
                                   aws_secret_access_key=aws_secret_key,
                                   endpoint_url=endpoint_url)
        else:
            self.s3 = boto3.client('s3')  # Use default credentials (e.g., IAM role)
        self.bucket_name = bucket_name

    def download_file(self, key):
        response = self.s3.get_object(Bucket=self.bucket_name, Key=key)
        return BytesIO(response['Body'].read())


    def upload_file(self, file_stream, key, content_type=None):
        if not content_type:
            content_type, _ = mimetypes.guess_type(key)
            if not content_type:
                content_type = 'application/octet-stream'  # fallback
        file_stream.seek(0)  # make sure the stream is at the beginning
        self.s3.put_object(
            Bucket=self.bucket_name,
            Key=key,
            Body=file_stream,
            ContentType=content_type
        )
        return f"http://46.101.215.10:9000/{self.bucket_name}/{key}"


    def extract_text(self, file_stream, file_name):
        ext = os.path.splitext(file_name)[1].lower()
        if ext == '.pdf':
            return self._extract_pdf(file_stream)
        elif ext == '.docx':
            return self._extract_docx(file_stream)
        elif ext == '.pptx':
            return self._extract_pptx(file_stream)
        elif ext == '.txt':
            return file_stream.read().decode('utf-8')
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

    def _extract_pdf(self, file_stream):
        reader = PyPDF2.PdfReader(file_stream)
        return " ".join(page.extract_text() for page in reader.pages if page.extract_text())

    def _extract_docx(self, file_stream):
    # Extract paragraphs outside of tables
        document = Document(file_stream)
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        
        # Extract table content
        table_texts = []
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        table_texts.append(cell_text)

        return " ".join(paragraphs + table_texts)

    def _extract_pptx(self, file_stream):
        prs = Presentation(file_stream)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return " ".join(text_runs)
    
    def get_file_text(self, key):
        file_stream = self.download_file(key)
        return self.extract_text(file_stream, key)

class TextSummarizer:
    def __init__(self):
        import openai
        self.openai = openai
        self.openai.api_key = 'OAI_KEY'
        self.model = 'gpt-3.5-turbo'

    def summarize(self, text):
        # Replace this with your actual LLM API
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",  # Or gpt-4, etc.
            messages=[
                {"role": "system", "content": "Create readable and informative conspect from the following document so student will understand everything, make it raw markdown:"},
                {"role": "user", "content": text}
            ],
            # max_tokens=max_tokens,
            temperature=0,
        )
        return response.choices[0].message.content